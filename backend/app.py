import os
import uuid
import threading
import time
from flask import Flask, request, jsonify, send_file
from werkzeug.utils import secure_filename
from flask_cors import CORS

import fitz  
from pptx import Presentation
from pptx.util import Inches, Pt

UPLOAD_FOLDER = './uploads'
CONVERTED_FOLDER = './converted'
ALLOWED_EXTENSIONS = {'pdf'}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER

CORS(app, origins=["http://localhost:3000"])  # adjust the origin as needed for your frontend

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

conversion_store = {}


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def convert_pdf_to_pptx(pdf_path, pptx_path):
    prs = Presentation()
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap()
        img_path = pptx_path + f'_page_{page_num}.png'
        pix.save(img_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)
    prs.save(pptx_path)


def background_conversion(upload_id, pdf_path, pptx_path):
    try:
        conversion_store[upload_id]['status'] = 'processing'
        time.sleep(2) 
        convert_pdf_to_pptx(pdf_path, pptx_path)
        conversion_store[upload_id]['status'] = 'completed'
        conversion_store[upload_id]['downloadUrl'] = f'/api/download/{upload_id}'
        conversion_store[upload_id]['filename'] = os.path.basename(pptx_path)
        conversion_store[upload_id]['fileSize'] = os.path.getsize(pptx_path)
    except Exception as e:
        conversion_store[upload_id]['status'] = 'error'
        conversion_store[upload_id]['error'] = str(e)


@app.route('/api/upload', methods=['POST'])
def upload_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        upload_id = str(uuid.uuid4())
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], f'{upload_id}.pdf')
        pptx_path = os.path.join(app.config['CONVERTED_FOLDER'], f'{upload_id}.pptx')
        file.save(pdf_path)
        conversion_store[upload_id] = {
            'status': 'processing',
            'progress': 0,
            'pdf_path': pdf_path,
            'pptx_path': pptx_path
        }
        
        threading.Thread(target=background_conversion, args=(upload_id, pdf_path, pptx_path)).start()
        return jsonify({'uploadId': upload_id})
    return jsonify({'error': 'Invalid file type'}), 400


@app.route('/api/status/<upload_id>', methods=['GET'])
def check_status(upload_id):
    info = conversion_store.get(upload_id)
    if not info:
        return jsonify({'error': 'Invalid uploadId'}), 404
    status = info['status']
    progress = 100 if status == 'completed' else (0 if status == 'processing' else 0)
    resp = {
        'status': status,
        'progress': progress,
        'uploadId': upload_id
    }
    if status == 'completed':
        resp['downloadUrl'] = f'/api/download/{upload_id}'
    if status == 'error':
        resp['error'] = info.get('error', 'Unknown error')
    return jsonify(resp)


@app.route('/api/result/<upload_id>', methods=['GET'])
def get_result(upload_id):
    info = conversion_store.get(upload_id)
    if not info or info['status'] != 'completed':
        return jsonify({'error': 'Not ready'}), 404
    return jsonify({
        'downloadUrl': f'/api/download/{upload_id}',
        'filename': info.get('filename', 'converted-presentation.pptx'),
        'fileSize': info.get('fileSize', 0)
    })


@app.route('/api/download/<upload_id>', methods=['GET'])
def download_file(upload_id):
    info = conversion_store.get(upload_id)
    if not info or info['status'] != 'completed':
        return jsonify({'error': 'Not ready'}), 404
    pptx_path = info['pptx_path']
    return send_file(pptx_path, as_attachment=True, download_name=info.get('filename', 'converted-presentation.pptx'))


@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8000, debug=True)
